from flask import Flask, request, jsonify
from flask_cors import CORS
from langchain_community.document_loaders import PyPDFLoader, DirectoryLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_ollama import OllamaEmbeddings
from langchain_ollama.chat_models import ChatOllama
from langchain.prompts import ChatPromptTemplate, PromptTemplate
from langchain.retrievers.multi_query import MultiQueryRetriever
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
import os
import shutil
from werkzeug.utils import secure_filename
from pptx import Presentation  # For PPTX files
from docx import Document  # For DOCX files

app = Flask(__name__)
CORS(app, origins=["http://localhost:3000"])  # Allow requests from your React frontend

# Configuration for file uploads
UPLOAD_FOLDER = "user_uploads"
ALLOWED_EXTENSIONS = {"pdf", "pptx", "docx"}  # Added support for PPT and DOCX
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
USER_INDEX_PATH = "user_faiss_index"
SYSTEM_INDEX_PATH = "system_faiss_index"
SYSTEM_DATA_PATH = "system_data"  # Folder containing system documents

# Ensure the upload folder exists
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(SYSTEM_DATA_PATH, exist_ok=True)  # Ensure system data folder exists

# Helper functions to load PPT and DOCX files
def load_pptx(file_path):
    """Load text from a PowerPoint (.pptx) file."""
    presentation = Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def load_docx(file_path):
    """Load text from a Word (.docx) file."""
    doc = Document(file_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

# Load and process documents
def load_and_process_documents(data_path, index_path="faiss_index", is_directory=True):
    try:
        if os.path.exists(index_path):
            # Load existing FAISS index
            vector_db = FAISS.load_local(index_path, OllamaEmbeddings(model="nomic-embed-text"), allow_dangerous_deserialization=True)
            print(f"Loaded existing FAISS index from {index_path}")
        else:
            documents = []
            if is_directory:
                # Load all files in the directory
                for filename in os.listdir(data_path):
                    file_path = os.path.join(data_path, filename)
                    if filename.endswith(".pdf"):
                        loader = PyPDFLoader(file_path)
                        documents.extend(loader.load())
                    elif filename.endswith(".pptx"):
                        text = load_pptx(file_path)
                        documents.append({"page_content": text, "metadata": {"source": file_path}})
                    elif filename.endswith(".docx"):
                        text = load_docx(file_path)
                        documents.append({"page_content": text, "metadata": {"source": file_path}})
            else:
                # Load a single file
                if data_path.endswith(".pdf"):
                    loader = PyPDFLoader(data_path)
                    documents = loader.load()
                elif data_path.endswith(".pptx"):
                    text = load_pptx(data_path)
                    documents.append({"page_content": text, "metadata": {"source": data_path}})
                elif data_path.endswith(".docx"):
                    text = load_docx(data_path)
                    documents.append({"page_content": text, "metadata": {"source": data_path}})
            
            print(f"Loaded {len(documents)} documents")
            
            # Split text into chunks
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
            chunks = text_splitter.split_documents(documents)
            print(f"Text split into {len(chunks)} chunks")
            
            # Create FAISS vector database
            vector_db = FAISS.from_documents(
                documents=chunks,
                embedding=OllamaEmbeddings(model="nomic-embed-text")
            )
            # Save the FAISS index
            vector_db.save_local(index_path)
            print(f"FAISS vector database created and saved at {index_path}")
        
        return vector_db
    except Exception as e:
        print(f"Error in load_and_process_documents: {e}")
        raise

# Set up LLM and retrieval
local_model = "llama3.2"  
llm = ChatOllama(model=local_model)

# Query prompt template
QUERY_PROMPT = PromptTemplate(
    input_variables=["question"],
    template="""You are an AI language model assistant. Your task is to generate 2
    different versions of the given user question to retrieve relevant documents from
    a vector database. By generating multiple perspectives on the user question, your
    goal is to help the user overcome some of the limitations of the distance-based
    similarity search. Provide these alternative questions separated by newlines.
    Original question: {question}""",
)

# Define the prompt template
template = """Some points to note while generating an answer:
a) Be courteous and creative in your responses. Answer like you are a friendly agent, making the user feel comfortable.
a.1) Imagine yourself as an intelligent assistant which is helping users with different types of study materials. You will have a knowledge base of study materials to search on.
b) Be very structured in your response. Provide the response with HTML tags as explained below:
- Use <b> tags to bold important words.
- Use <br> tags to display contents in new lines wherever required for clear displaying.
- Use <ul> and <li> tags for bullet points.
- Use <p> tags for paragraphs.
Important Note: We need to display the content with proper HTML tags. Your task is to automatically add the relevant <li>, <ul>, <br>, <p>, and <b> tags to beautifully present the responses to be displayed in an HTML <div> tag.
c) The answer should not be in the email format and should look like a normal chat.
d) **Important**: Do not use Markdown formatting (e.g., `**bold**`, `*italic*`, `- bullet points`). Use only HTML tags for formatting.

Context:
{context}

Question: {question}

Answer:
Please provide the answer in a clear, point-wise format for better readability using HTML tags:
<p>1. ...</p>
<p>2. ...</p>
<p>3. ...</p>
"""
prompt = ChatPromptTemplate.from_template(template)
print("Prompt work successfully")

# Global retriever and chain for both user and system
user_vector_db = None
user_retriever = None
user_chain = None

system_vector_db = None
system_retriever = None
system_chain = None

def setup_system_chain():
    """Initialize the system RAG chain with default documents"""
    global system_vector_db, system_retriever, system_chain
    
    try:
        # Load system documents (from SYSTEM_DATA_PATH directory)
        system_vector_db = load_and_process_documents(SYSTEM_DATA_PATH, index_path=SYSTEM_INDEX_PATH, is_directory=True)
        
        # Set up retriever
        system_retriever = MultiQueryRetriever.from_llm(
            system_vector_db.as_retriever(), 
            llm,
            prompt=QUERY_PROMPT
        )
        
        # Create chain
        system_chain = (
            {"context": system_retriever, "question": RunnablePassthrough()}
            | prompt
            | llm
            | StrOutputParser()
        )
        
        print("System chain initialized successfully")
        return True
    except Exception as e:
        print(f"Error setting up system chain: {e}")
        return False

def setup_user_chain(file_path):
    """Initialize the user RAG chain with uploaded document"""
    global user_vector_db, user_retriever, user_chain
    
    try:
        # Process the uploaded file (single file, not a directory)
        user_vector_db = load_and_process_documents(file_path, index_path=USER_INDEX_PATH, is_directory=False)
        
        # Set up retriever
        user_retriever = MultiQueryRetriever.from_llm(
            user_vector_db.as_retriever(), 
            llm,
            prompt=QUERY_PROMPT
        )
        
        # Create chain
        user_chain = (
            {"context": user_retriever, "question": RunnablePassthrough()}
            | prompt
            | llm
            | StrOutputParser()
        )
        
        print("User chain initialized successfully")
        return True
    except Exception as e:
        print(f"Error setting up user chain: {e}")
        return False

def delete_user_data():
    """Delete user data and reset the user chain"""
    global user_vector_db, user_retriever, user_chain
    
    # Delete the FAISS index files
    if os.path.exists(USER_INDEX_PATH):
        try:
            shutil.rmtree(USER_INDEX_PATH)
            print(f"Deleted user FAISS index at {USER_INDEX_PATH}")
        except Exception as e:
            print(f"Error deleting user FAISS index: {e}")
    
    # Delete uploaded files
    if os.path.exists(UPLOAD_FOLDER):
        try:
            shutil.rmtree(UPLOAD_FOLDER)
            os.makedirs(UPLOAD_FOLDER, exist_ok=True)
            print(f"Cleared upload folder at {UPLOAD_FOLDER}")
        except Exception as e:
            print(f"Error clearing upload folder: {e}")
    
    # Reset the global variables
    user_vector_db = None
    user_retriever = None
    user_chain = None

# Initialize system chain at startup
setup_system_chain()

# Root route
@app.route('/')
def home():
    return "Welcome to the PDF Chatbot API! Use the /upload endpoint to upload a file and the /chat endpoint to interact with the chatbot."

# File upload route
@app.route('/upload', methods=['POST'])
def upload_file():
    """
    Upload a file (PDF, PPT, or DOCX) and process it for the chatbot.
    """
    if "file" not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No selected file"}), 400
    
    if file and file.filename.endswith(tuple(ALLOWED_EXTENSIONS)):
        # First, delete any existing user data
        delete_user_data()
        
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
        file.save(file_path)
        
        if setup_user_chain(file_path):
            return jsonify({
                "message": "File uploaded and processed successfully",
                "using_user_file": True
            }), 200
        else:
            return jsonify({"error": "Failed to process the file"}), 500
    else:
        return jsonify({"error": "Invalid file type. Only PDF, PPT, and DOCX files are allowed"}), 400

# Delete user data route
@app.route('/delete-user-data', methods=['POST'])
def delete_user_data_endpoint():
    """
    Delete the user's uploaded files and FAISS index.
    """
    try:
        delete_user_data()
        return jsonify({
            "message": "User data deleted successfully",
            "using_user_file": False
        }), 200
    except Exception as e:
        return jsonify({"error": f"Failed to delete user data: {e}"}), 500

# Chat route
@app.route('/chat', methods=['POST'])
def chat_with_file():
    """
    Chat with either the uploaded file or system documents using the RAG chain.
    """
    data = request.json
    question = data.get('question')
    
    if not question:
        return jsonify({"error": "No question provided"}), 400
    
    # Determine which chain to use (user chain takes precedence)
    active_chain = user_chain if user_chain else system_chain
    using_user_file = user_chain is not None
    
    if not active_chain:
        return jsonify({"error": "No documents available. Please upload a file or ensure system documents exist."}), 400
    
    try:
        response = active_chain.invoke(question)
        return jsonify({
            "response": response,
            "using_user_file": using_user_file
        })
    except Exception as e:
        print(f"Error in chain invocation: {e}")
        return jsonify({"error": "Failed to generate response"}), 500

if __name__ == '__main__':
    app.run(debug=True, port=8080)